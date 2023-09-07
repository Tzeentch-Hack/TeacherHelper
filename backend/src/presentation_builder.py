from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx import Presentation
from PIL import Image, ImageEnhance
import io
import subprocess
import os


# Constants
DEFAULT_FONT = 'Palatino'
TITLE_FONT_SIZE = Pt(48)
TEXT_FONT_SIZE = Pt(24)
TITLE_COLOR = RGBColor(235, 215, 215)  # White color text
SIZE_DECREMENT = Pt(1)
MAX_TITLE_SIZE = Pt(48)
MIN_TITLE_SIZE = Pt(20)
MAX_TEXT_SIZE = Pt(24)
MIN_TEXT_SIZE = Pt(12)

def fit_font_to_text(text, font, max_width, max_size, min_size):
    estimated_width = len(text) * font.size.pt * 7200
    while estimated_width > max_width and font.size > min_size:
        font.size -= SIZE_DECREMENT
        estimated_width = len(text) * font.size.pt * 7200

    if font.size > max_size:
        font.size = max_size


def add_custom_title(slide, prs, title):
    slide_width = prs.slide_width
    left = top = Inches(0)
    width = slide_width
    height = Inches(1)
    text = title
    txBox = add_textbox(slide, left, top, width, height, text, DEFAULT_FONT, TITLE_FONT_SIZE, TITLE_COLOR, MAX_TITLE_SIZE, MIN_TITLE_SIZE)
    txBox.left = (prs.slide_width - txBox.width) // 2

def add_textbox(slide, left, top, width, height, text, font_name, font_size, font_color, max_size, min_size):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = text
    run = p.runs[0]
    font = run.font
    font.name = font_name
    font.size = font_size
    font.color.rgb = font_color

    fit_font_to_text(text, font, width, max_size, min_size)
    return txBox

def convert_pptx_to_pdf(input_path: str, output_folder: str = None) -> None:
    cmd = [
        "soffice",
        "--headless",
        "--convert-to", "pdf"
    ]

    if output_folder:
        cmd.extend(["--outdir", output_folder])

    cmd.append(input_path)

    try:
        subprocess.run(cmd, check=True)
        print(f"File {input_path} converted successfully!")
    except subprocess.CalledProcessError:
        print(f"Error converting {input_path}.")
        raise




def convert_pdf_to_jpeg(pdf_path: str, output_name: str, output_folder: str = None, page: int = 1, resolution: int = 120, quality: int = 90) -> None:
    if output_folder:
        output_prefix = os.path.join(output_folder, output_name)
    else:
        output_prefix = output_name

    cmd = [
        "pdftoppm",
        "-singlefile",
        "-f", str(page),
        "-r", str(resolution),
        "-jpeg",
        "-jpegopt", f"quality={quality}",
        pdf_path,
        output_prefix
    ]

    try:
        subprocess.run(cmd, check=True)
        print(f"File {pdf_path} converted successfully!")
    except subprocess.CalledProcessError:
        print(f"Error converting {pdf_path}.")
        raise


def convert_pptx_to_jpg(input_path: str, output_name: str, output_folder: str = None, page: int = 1,
                        resolution: int = 120, quality: int = 90) -> None:
    # Step 1: Convert pptx to pdf
    pdf_output_folder = None
    if output_folder:
        pdf_output_folder = os.path.join(output_folder, "temp")
        os.makedirs(pdf_output_folder, exist_ok=True)

    pdf_name = os.path.splitext(os.path.basename(input_path))[0] + '.pdf'
    pdf_output_path = os.path.join(pdf_output_folder if pdf_output_folder else '', pdf_name)

    convert_pptx_to_pdf(input_path, pdf_output_folder)

    convert_pdf_to_jpeg(pdf_output_path, output_name, output_folder, page, resolution, quality)

    os.remove(pdf_output_path)



def add_text(slide, prs, phrases):
    slide_width = prs.slide_width
    left = 0
    top = Inches(1.5)
    width = slide_width
    height = Inches(2)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.margin_top = Inches(0.5)
    tf.text_anchor = 'middle'

    for line in phrases:
        p = tf.add_paragraph()
        p.text = line
        p.space_after = Inches(0.4)
        run = p.runs[0]
        font = run.font
        font.name = DEFAULT_FONT
        font.size = TEXT_FONT_SIZE
        font.color.rgb = TITLE_COLOR

        fit_font_to_text(line, font, width, MAX_TEXT_SIZE, MIN_TEXT_SIZE)

    txBox.left = int((prs.slide_width - txBox.width) // 2)

def darken_image(img_path, factor=0.5):
    with Image.open(img_path) as im:
        enhancer = ImageEnhance.Brightness(im)
        im = enhancer.enhance(factor)
        byte_io = io.BytesIO()
        im.save(byte_io, format='PNG')
        byte_io.seek(0)
    return byte_io

def create_pptx_and_convert_to_jpg(title: str, phrases: list, img_path: str, output_path: str) -> str:
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    darkened_img_stream = darken_image(img_path)

    left = top = 0
    slide.shapes.add_picture(darkened_img_stream, left, top, width=prs.slide_width, height=prs.slide_height)

    add_custom_title(slide, prs, title)
    add_text(slide, prs, phrases)

    temp_pptx_path = "temp_presentation.pptx"
    prs.save(temp_pptx_path)

    output_name = os.path.basename(output_path)
    output_folder = os.path.dirname(output_path)
    convert_pptx_to_jpg(temp_pptx_path, output_name, output_folder)
    os.remove(temp_pptx_path)

    return output_path


def create_pptx_multiple_slides(slides_data, output_path):
    prs = Presentation()
    for slide_item in slides_data:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        darkened_img_stream = darken_image(slide_item["img_path"])

        left = top = 0
        slide.shapes.add_picture(darkened_img_stream, left, top, width=prs.slide_width, height=prs.slide_height)

        add_custom_title(slide, prs, slide_item["topic_name"])
        add_text(slide, prs, slide_item["options"])

    temp_pptx_path = output_path
    prs.save(temp_pptx_path)


if __name__ == "__main__":
    slides_data = [
        {
            "topic_name": "O'zbekiston Vatanim manim",
            "options": ["Домла педарас, мраз", "Igraka ohuenniy builder presentatsiy  ohuenniy builder builder builder builder", "Sanya sexi"],
            "img_path": "../data/pictures/temp/2.png"
        },
        {
            "topic_name": "Second Slide Title",
            "options": ["Phrase 1", "Phrase 2", "Phrase 3"],
            "img_path": "../data/pictures/temp/1.png"
        }
    ]
    create_pptx_multiple_slides(slides_data, "./multiple.pptx")